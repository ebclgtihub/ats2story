#!/usr/bin/env python3
"""
ats2pptx — Proof-of-Concept Konverter: imc Content Studio (.ats) -> PowerPoint (.pptx)

Eine .ats-Datei ist ein ZIP mit:
  document/document.xml   -> Kursstruktur (Ordner=Kapitel, animation=Folie, exam=Test)
  resources/ata/*.ata     -> je Folie ein verschachteltes ZIP mit eigenem document.xml + Bildern
  resources/ati/*.ati     -> je Interaktion ein verschachteltes ZIP (Drag&Drop / Quiz)

Dieser PoC überträgt die FOLIEN (Text + Bilder + Layout) nach PowerPoint.
Quizze/Interaktionen werden als TODO-Platzhalter eingefügt (separater Schritt).

Aufruf:
  python3 ats2pptx.py <input.ats> <output.pptx> [--max N] [--chapters "A,B"]
"""
import sys, os, zipfile, io, html, re, argparse
import xml.etree.ElementTree as ET
from html.parser import HTMLParser

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NS = '{http://im-c.de/xml/authoring/1.0}'
EMU_PER_PX = 9525            # 96 DPI -> 1 px = 9525 EMU (Koordinaten 1:1 übernehmen)
CANVAS_W, CANVAS_H = 1024, 748

def loc(tag): return tag.split('}')[-1]

# ---------- Rich-Text Parser (HTML im richText-Attribut) ----------
class RichTextParser(HTMLParser):
    """Wandelt das HTML aus richText in [paragraph][run(text, bold, italic, underline, color)] um."""
    def __init__(self):
        super().__init__()
        self.paras = [[]]
        self.bold = self.italic = self.underline = False
        self.color = None
        self._style_stack = []
    def handle_starttag(self, tag, attrs):
        d = dict(attrs)
        if tag == 'p':
            if self.paras[-1]:
                self.paras.append([])
        elif tag == 'br':
            self.paras.append([])
        elif tag in ('span', 'font', 'b', 'strong', 'i', 'em', 'u'):
            st = d.get('style', '')
            b = self.bold or tag in ('b', 'strong') or 'font-weight: bold' in st or 'font-weight:bold' in st
            it = self.italic or tag in ('i', 'em') or 'font-style: italic' in st or 'font-style:italic' in st
            un = self.underline or tag == 'u' or 'underline' in st
            col = self.color
            m = re.search(r'color:\s*(#[0-9a-fA-F]{6})', st)
            if m: col = m.group(1)
            self._style_stack.append((self.bold, self.italic, self.underline, self.color))
            self.bold, self.italic, self.underline, self.color = b, it, un, col
    def handle_endtag(self, tag):
        if tag in ('span', 'font', 'b', 'strong', 'i', 'em', 'u') and self._style_stack:
            self.bold, self.italic, self.underline, self.color = self._style_stack.pop()
    def handle_data(self, data):
        if data:
            self.paras[-1].append((data, self.bold, self.italic, self.underline, self.color))
    def get(self):
        return [p for p in self.paras if p] or [[]]

def parse_rich(rich):
    if not rich: return [[]]
    p = RichTextParser()
    p.feed(html.unescape(rich))
    return p.get()

# ---------- Geometrie ----------
def rect_of(el):
    r = el.find(NS+'complexproperty/'+NS+'rect')
    if r is None:
        r = el.find('.//'+NS+'rect')
    if r is None: return None
    return (int(float(r.get('x',0))), int(float(r.get('y',0))),
            int(float(r.get('width',100))), int(float(r.get('height',30))))

def emu(px): return Emu(int(px*EMU_PER_PX))

# ---------- .ata Folie -> PPTX-Folie ----------
def build_slide(prs, ata_bytes, title, chapter):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    notes = slide.notes_slide.notes_text_frame
    notes.text = f"Kapitel: {chapter}\nFolie: {title}"

    z = zipfile.ZipFile(io.BytesIO(ata_bytes))
    try:
        doc = ET.fromstring(z.read('document/document.xml'))
    except KeyError:
        return slide
    # Elemente in Layer-Reihenfolge: erst Bilder, dann Text (grob nach 'layer')
    elems = []
    for el in doc:
        lt = loc(el.tag)
        if lt in ('image', 'text', 'rectangle'):
            elems.append((int(el.get('layer', 0)), lt, el))
    elems.sort(key=lambda e: e[0])

    for _, lt, el in elems:
        rect = rect_of(el)
        if rect is None: continue
        x, y, w, h = rect
        if lt == 'image':
            res = el.find('.//'+NS+'resource')
            if res is None: continue
            path = res.get('path')
            try:
                img = z.read(path)
            except KeyError:
                continue
            try:
                slide.shapes.add_picture(io.BytesIO(img), emu(x), emu(y), emu(w), emu(h))
            except Exception:
                pass
        elif lt == 'text':
            paras = parse_rich(el.get('richText', ''))
            if not any(any(r[0].strip() for r in p) for p in paras):
                continue
            tb = slide.shapes.add_textbox(emu(x), emu(y), emu(max(w, 30)), emu(max(h, 20)))
            tf = tb.text_frame; tf.word_wrap = True
            base_size = int(float(el.get('fontSize', 14)))
            base_font = el.get('fontFamily', 'Arial')
            base_color = el.get('textColor', '#000000')
            align = el.get('textAlign', 'left')
            for i, para in enumerate(paras):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
                               'justify': PP_ALIGN.JUSTIFY}.get(align, PP_ALIGN.LEFT)
                for (text, b, it, un, col) in para:
                    run = p.add_run(); run.text = text
                    f = run.font
                    f.size = Pt(base_size); f.name = base_font
                    f.bold = b; f.italic = it; f.underline = un
                    c = col or base_color
                    try: f.color.rgb = RGBColor.from_string(c.lstrip('#'))
                    except Exception: pass
    return slide

# ---------- Hauptlauf ----------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('inp'); ap.add_argument('out')
    ap.add_argument('--max', type=int, default=0, help='max. Folien (0=alle)')
    ap.add_argument('--chapters', default='', help='nur diese Kapitel (Komma-getrennt, Teilstring)')
    a = ap.parse_args()

    ats = zipfile.ZipFile(a.inp)
    root = ET.fromstring(ats.read('document/document.xml'))
    chap_filter = [c.strip().lower() for c in a.chapters.split(',') if c.strip()]

    prs = Presentation()
    prs.slide_width = emu(CANVAS_W); prs.slide_height = emu(CANVAS_H)

    made = exams = 0
    def walk(e, chapter):
        nonlocal made, exams
        lt = loc(e.tag)
        if lt == 'folder':
            chapter = e.get('name') or e.get('displayName') or chapter
        if lt == 'animation':
            if a.max and made >= a.max: return
            if chap_filter and not any(cf in (chapter or '').lower() for cf in chap_filter):
                pass
            else:
                res = e.find('.//'+NS+'resource')
                name = e.get('name') or '(ohne Titel)'
                if res is not None and res.get('path'):
                    try:
                        ata = ats.read(res.get('path'))
                        build_slide(prs, ata, name, chapter or '')
                        made += 1
                        if made % 20 == 0: print(f'  ... {made} Folien')
                    except KeyError:
                        pass
            return
        if lt == 'exam':
            if not chap_filter or any(cf in (chapter or '').lower() for cf in chap_filter):
                blank = prs.slide_layouts[6]
                s = prs.slides.add_slide(blank)
                tb = s.shapes.add_textbox(emu(80), emu(300), emu(860), emu(150))
                tf = tb.text_frame; tf.word_wrap = True
                tf.text = f'[TEST] {e.get("name","")}'
                tf.paragraphs[0].runs[0].font.size = Pt(28)
                tf.paragraphs[0].runs[0].font.bold = True
                p = tf.add_paragraph(); p.text = 'Quiz-Fragen separat nach Storyline übertragen (siehe Report).'
                p.runs[0].font.size = Pt(14)
                exams += 1
        for c in list(e):
            if a.max and made >= a.max and lt != 'document':
                break
            walk(c, chapter)
    walk(root, None)

    prs.save(a.out)
    print(f'\nFertig: {made} Folien + {exams} Test-Platzhalter -> {a.out}')

if __name__ == '__main__':
    main()
